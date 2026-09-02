# -*- coding: utf-8 -*-
"""生成认证交付物《项目介绍.pptx》（docs/交付文档/项目介绍.pptx）。

仅本地生成用，依赖 python-pptx（已在本地 venv 安装）。
运行：python tools/make_ppt.py

逐页文案见 docs/交付文档/PPT文案.md。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "docs", "交付文档", "项目介绍.pptx")
WEB_IMG = os.path.join(ROOT, "web", "images")

# 配色（新中式墨绿 + 米白 + 朱红点缀）
INK = RGBColor(0x1F, 0x3D, 0x2E)      # 墨绿
RICE = RGBColor(0xF5, 0xF1, 0xE6)     # 米白
VERM = RGBColor(0xB5, 0x3A, 0x29)     # 朱红
GOLD = RGBColor(0xC8, 0xA9, 0x5C)     # 金
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GREY = RGBColor(0x6B, 0x6B, 0x6B)
PANEL = RGBColor(0xEC, 0xE7, 0xD8)    # 浅米面板

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_bg(slide, color=RICE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def set_text(tf, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tf = tf.text_frame if hasattr(tf, "text_frame") else tf
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def add_para(tf, text, size=16, color=DARK, bold=False, align=PP_ALIGN.LEFT,
             level=0, bullet=False, space=6, font="Microsoft YaHei"):
    tf = tf.text_frame if hasattr(tf, "text_frame") else tf
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = ("• " + text) if bullet else text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def title_bar(slide, text, sub=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.12)
    set_text(tf, text, size=28, color=RICE, bold=True)
    if sub:
        add_para(tf, sub, size=13, color=GOLD, space=0)


def sec(slide, l, t, w, h, text, color=INK):
    """小标题分区条。"""
    bar = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, text, size=16, color=RICE, bold=True)


def pic_fit(slide, path, l, t, w, h):
    """等比缩放嵌入图片到指定框内并居中。"""
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
    except Exception:
        return
    box_ar = w / h
    img_ar = iw / ih
    if img_ar > box_ar:
        nw = w; nh = w / img_ar
    else:
        nh = h; nw = h * img_ar
    nl = l + (w - nw) / 2
    nt = t + (h - nh) / 2
    slide.shapes.add_picture(path, Inches(nl), Inches(nt), Inches(nw), Inches(nh))


def add_table(slide, l, t, w, h, header, data, col_w=None, fsize=11):
    nrow = len(data) + 1
    ncol = len(header)
    gf = slide.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = True
    if col_w:
        for c, cw in enumerate(col_w):
            tbl.columns[c].width = Inches(cw)
    for c in range(ncol):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        set_text(cell, header[c], size=fsize + 1, color=RICE, bold=True)
    for r in range(1, nrow):
        for c in range(ncol):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RICE if r % 2 else PANEL
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            set_text(cell, data[r - 1][c], size=fsize, color=DARK)
    return tbl


def footnote(slide, text="免责声明：本系统为中医药科普辅助工具，不构成医疗建议。"):
    tb = box(slide, 0.5, 7.05, 12.3, 0.35)
    set_text(tb, text, size=10, color=GREY, align=PP_ALIGN.CENTER)


def show(fn):
    return os.path.join(WEB_IMG, "showcase", fn) if os.path.exists(
        os.path.join(WEB_IMG, "showcase", fn)) else None


# ---------------- 1. 封面 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s, INK)
plum = os.path.join(WEB_IMG, "abstract-plum.png")
if os.path.exists(plum):
    s.shapes.add_picture(plum, Inches(9.6), Inches(4.6), Inches(3.5), Inches(2.8))
tf = box(s, 1.0, 2.1, 11.3, 2.8); tf.text_frame.word_wrap = True
set_text(tf, "中草药多模态识别系统", size=46, color=RICE, bold=True)
add_para(tf, "本草掠影 · Herb Recognition", size=22, color=GOLD, space=4)
add_para(tf, "视觉 + 文本 + 知识图谱融合的中医药智能识别与药性分析平台", size=16, color=RICE, space=0)
tf2 = box(s, 1.0, 6.4, 11.3, 0.8)
set_text(tf2, "拍一张药材照片，即可识别、解析药性、推演配伍与方剂 — 认证式交付物 · 项目介绍", size=13, color=GOLD)

# ---------------- 2. 目录 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "目录", "Contents")
toc = [
    "一、项目背景与痛点", "八、功能总览（一表看懂）",
    "二、系统总览与总体架构", "九、官方评估指标",
    "三、模块一 · 多模态识别引擎", "十、技术栈与部署",
    "四、模块二 · 中医药知识图谱", "十一、创新点",
    "五、模块三 · 特性检索与方剂推荐", "十二、不足与改进方向",
    "六、模块四 · RAG 问答与可解释性", "十三、总结与交付物",
    "七、模块五 · 交互展示与平台接入",
]
lt = box(s, 0.7, 1.5, 6.0, 5.3); lt.text_frame.word_wrap = True
rt = box(s, 6.9, 1.5, 6.0, 5.3); rt.text_frame.word_wrap = True
for i, item in enumerate(toc):
    tgt = lt if i < 7 else rt
    p = tgt.text_frame.paragraphs[0] if i in (0, 7) else tgt.text_frame.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run(); r.text = item
    r.font.size = Pt(17); r.font.color.rgb = INK if i % 2 == 0 else DARK
    r.font.bold = (i % 2 == 0); r.font.name = "Microsoft YaHei"

# ---------------- 3. 项目背景与痛点 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "一、项目背景与痛点", "Why — 普通人识别与理解中草药的现实困难")
tf = box(s, 0.6, 1.45, 12.1, 5.4); tf.text_frame.word_wrap = True
set_text(tf, "为什么做？", size=18, color=INK, bold=True)
add_para(tf, "种类繁多、外观相似度高：传统辨认高度依赖经验，普通用户/基层人员难快速准确识别。", size=15, bullet=True)
add_para(tf, "药性专业且零散：性味、归经、配伍禁忌、方剂知识体量大，散落古籍药典，查阅效率低。", size=15, bullet=True)
add_para(tf, "工具割裂：市面识别工具大多「只识图、不识药」——识别完即结束，缺中医药知识联动。", size=15, bullet=True)
add_para(tf, "共性风险：相似药材极易混淆，误识别+误用药=健康风险；缺置信度提示与风险警示。", size=15, color=VERM, bold=True, bullet=True)
add_para(tf, "目标：打通「识别—解析—推荐—问答」全链路，并标注知识来源以抑制大模型幻觉；定位科普辅助，全程医疗风险提示。", size=15, color=INK, bold=True, bullet=True)
footnote(s)

# ---------------- 4. 系统总览与总体架构 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "二、系统总览与总体架构", "How — 四层架构：输入 → 多模态引擎 → 知识服务 → 交互")
tf = box(s, 0.6, 1.45, 12.1, 5.4); tf.text_frame.word_wrap = True
set_text(tf, "输入层", size=16, color=INK, bold=True)
add_para(tf, "药材图片（原植物/饮片/干药材）+ 可选文本（性味·归经·功效·提问）；支持拖拽/点击/Ctrl+V。", size=14, bullet=True)
add_para(tf, "多模态识别引擎（核心）", size=16, color=INK, bold=True)
add_para(tf, "Swin-Tiny 视觉编码 + BERT 中文编码 + HCA 跨模态注意力融合；双分支自动切换（有文本走多模态，无文本走纯视觉）。", size=14, bullet=True)
add_para(tf, "知识服务层", size=16, color=INK, bold=True)
add_para(tf, "中医药知识图谱（200+ 药材、12 首方剂、十八反/十九畏规则）+ RAG 语义检索 + 方剂/相似药推荐。", size=14, bullet=True)
add_para(tf, "交互层", size=16, color=INK, bold=True)
add_para(tf, "新中式 Web 单页应用 + Gradio 调试端 + FastAPI REST API（可服务化）。", size=14, color=VERM, bullet=True)
footnote(s)

# ---------------- 5. 模块一：多模态识别引擎（技术核心） ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "三、模块一 · 多模态识别引擎", "技术核心：视觉 + 文本 + HCA 融合，面向 8G 显存优化")
p = show("gancao.jpg")
if p:
    pic_fit(s, p, 0.6, 1.5, 5.2, 4.8)
tf = box(s, 6.2, 1.55, 6.6, 5.3); tf.text_frame.word_wrap = True
set_text(tf, "技术选型", size=17, color=INK, bold=True)
add_para(tf, "视觉：Swin-Tiny（window attention，约 2.6G；备选 EfficientNet-B0/ResNet50/ConvNeXt-Tiny）。", size=13.5, bullet=True)
add_para(tf, "文本：BERT-base-chinese（冻结主干，仅训投影层省算力）。", size=13.5, bullet=True)
add_para(tf, "融合：HCA-Fusion 层级式跨模态注意力（自研）。", size=13.5, bullet=True)
add_para(tf, "分类头：多模态分支 + 纯视觉保底双分支。", size=13.5, bullet=True)
set_text(tf, "关键设计", size=17, color=INK, bold=True)
add_para(tf, "模拟中医「望闻问切」互补：视觉/文本 token 互为 Q/K/V 交叉注意力，再自注意力精炼。", size=13.5, bullet=True)
add_para(tf, "双分支保底：文本为空自动走纯视觉，根本保证「无文本也能靠图识别」。", size=13.5, bullet=True)
add_para(tf, "启用 ImageNet 预训练微调：从零训练仅 15-20%，微调后纯视觉达 90%+。", size=13.5, color=VERM, bullet=True)
footnote(s)

# ---------------- 6. 模块一续：识别流程与可解释 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "模块一（续）· 推理流程与易混鉴别", "从像素到药材档案，并对相似药主动提示")
lt = box(s, 0.6, 1.5, 6.0, 5.2); lt.text_frame.word_wrap = True
set_text(lt, "推理流程", size=17, color=INK, bold=True)
add_para(lt, "图片 → Swin-Tiny 提特征 → 投影到 embed_dim(512)。", size=14, bullet=True)
add_para(lt, "文本 → BERT 提 [CLS] → 投影到同维度。", size=14, bullet=True)
add_para(lt, "HCA 融合 → 多模态 logits（有文本）/ 纯视觉 logits（无文本）。", size=14, bullet=True)
add_para(lt, "输出 Top-5 识别 + 置信度徽章。", size=14, bullet=True)
rt = box(s, 6.9, 1.5, 6.0, 5.2); rt.text_frame.word_wrap = True
set_text(rt, "易混鉴别（16 组对照库）", size=17, color=INK, bold=True)
add_para(rt, "白术/苍术、半夏/天南星、人参/党参、枸杞/地骨皮…", size=14, bullet=True)
add_para(rt, "输出外观差异 + 简易鉴别法（嗅闻/口尝/水试）。", size=14, bullet=True)
add_para(rt, "连带功效与毒性提示，辅助人工复核，降低误用风险。", size=14, color=VERM, bullet=True)
footnote(s)

# ---------------- 7. 模块二：中医药知识图谱 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "四、模块二 · 中医药知识图谱", "数据驱动的关系网络 + 硬编码领域安全红线")
p = show("gouqizi.jpg")
if p:
    pic_fit(s, p, 0.6, 1.5, 5.2, 4.8)
tf = box(s, 6.2, 1.55, 6.6, 5.3); tf.text_frame.word_wrap = True
set_text(tf, "构建与规模", size=17, color=INK, bold=True)
add_para(tf, "内存版 NetworkX（可选 Neo4j 持久化），CSV 加载节点与多类型关系。", size=13.5, bullet=True)
add_para(tf, "200+ 药材节点、12 首经典方剂、151 条增量字段（别名/病症/个体禁忌）。", size=13.5, bullet=True)
set_text(tf, "关系类型（边）", size=17, color=INK, bold=True)
add_para(tf, "paired 相须相使 · incompatible 十八反 · restraint 十九畏 · category/meridian 功效与归经维度。", size=13.5, bullet=True)
set_text(tf, "内置安全红线", size=17, color=INK, bold=True)
add_para(tf, "硬编码十八反/十九畏歌诀，未收录药材也能被查询与可视化，避免盲区。", size=13.5, bullet=True)
add_para(tf, "毒性自动分级（大毒/有毒/小毒/微毒），强制弹窗警示。", size=13.5, color=VERM, bullet=True)
footnote(s)

# ---------------- 8. 模块三：特性检索与方剂推荐 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "五、模块三 · 特性检索与方剂推荐", "把识别结果自动联动药性，并规避危险配伍")
tf = box(s, 0.6, 1.5, 12.1, 5.3); tf.text_frame.word_wrap = True
set_text(tf, "特性检索", size=17, color=INK, bold=True)
add_para(tf, "自然语言输入「味甘平，归肝肾经」「清热明目」→ 解析为性味/归经/功效三类条件逐味打分。", size=14, bullet=True)
add_para(tf, "输出「完全匹配/部分匹配」卡片（印章式角标），支持按药材名批量查询。", size=14, bullet=True)
set_text(tf, "相似药推荐", size=17, color=INK, bold=True)
add_para(tf, "功效分类×3 + 功效 n-gram×2 + 性味×1 + 归经×1 综合打分，仅推功效相近者，辅助辨证选药。", size=14, bullet=True)
set_text(tf, "方剂推荐", size=17, color=INK, bold=True)
add_para(tf, "基于「主治功效匹配 + 常用配伍 + 禁忌规避」打分排序；命中十八反/十九畏直接剔除，从源头规避危险配伍。", size=14, bullet=True)
add_para(tf, "返回组成、来源、主治、用法、警示，并标注配伍风险。", size=14, color=VERM, bullet=True)
footnote(s)

# ---------------- 9. 模块四：RAG 问答与可解释性 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "六、模块四 · RAG 问答与可解释性", "可溯源回答 + 模型关注可视化，降低黑箱感")
p = show("chenpi.jpg")
if p:
    pic_fit(s, p, 0.6, 1.5, 5.2, 4.8)
tf = box(s, 6.2, 1.55, 6.6, 5.3); tf.text_frame.word_wrap = True
set_text(tf, "RAG 增强问答", size=17, color=INK, bold=True)
add_para(tf, "本地 BERT 对知识切片做中文语义向量检索，为 LLM 组装专家上下文。", size=13.5, bullet=True)
add_para(tf, "策略：药名精确命中优先 + 语义 Top-K 补充；回答附「知识库来源」折叠卡，减幻觉。", size=13.5, bullet=True)
add_para(tf, "离线可用、零新增依赖；BERT 失败自动降级关键词匹配，不抛异常。", size=13.5, bullet=True)
add_para(tf, "外部 LLM（智谱 GLM）无 Key/失败 → 自动降级返回本地图谱结果。", size=13.5, bullet=True)
set_text(tf, "可解释性", size=17, color=INK, bold=True)
add_para(tf, "Grad-CAM 热力图：滑块实时调透明度，展示模型「看哪里」。", size=13.5, bullet=True)
add_para(tf, "推理路径：识别→图谱查询→推荐 全链路可追溯，回答附来源标注。", size=13.5, color=VERM, bullet=True)
footnote(s)

# ---------------- 10. 模块五：交互展示与平台接入 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "七、模块五 · 交互展示与平台接入", "三端入口：新中式 Web / Gradio / REST API")
tf = box(s, 0.6, 1.5, 12.1, 5.3); tf.text_frame.word_wrap = True
set_text(tf, "新中式 Web 前端（推荐入口，原生 HTML/CSS/JS，FastAPI 托管）", size=16, color=INK, bold=True)
add_para(tf, "图片识别 / 特性检索 / Grad-CAM / AI 对话 / 药材关系图谱 五页签。", size=14, bullet=True)
add_para(tf, "关系图谱力导向网络（拖拽/缩放/点击）：红虚线=十八反、橙虚线=十九畏、绿实线=相须相使。", size=14, bullet=True)
set_text(tf, "Gradio 演示 / REST 服务化", size=16, color=INK, bold=True)
add_para(tf, "Gradio 调试端 5 页签与前端一致；python main.py --mode serve 提供 API。", size=14, bullet=True)
add_para(tf, "/predict 识别 · /search 检索 · /explain 热图 · /chat 对话 · /graph 图谱 · /herbs 列表。", size=14, bullet=True)
add_para(tf, "端侧：导出纯视觉分支为 TorchScript（tools/export_model.py），输出 [B,163] logits，无需文本模型。", size=14, color=VERM, bullet=True)
footnote(s)

# ---------------- 11. 功能总览（一表看懂） ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "八、功能总览", "一表看懂：输入 / 输出 / 技术")
data = [
    ["多模态识别", "图 + 文本", "Top-5 + 置信度", "Swin + BERT + HCA"],
    ["易混鉴别", "识别结果", "外观差异 + 鉴别法", "16 组对照库"],
    ["特性检索", "性味/归经/功效", "完全/部分匹配卡片", "条件解析 + 打分"],
    ["相似药推荐", "药材名", "功效相近药材", "多维度相似度"],
    ["方剂推荐", "主药 + 症状", "经典方剂 + 风险提示", "功效匹配 + 禁忌规避"],
    ["AI 对话", "文 + 图", "专家级解释 + 来源", "RAG + LLM"],
    ["关系图谱", "聚焦药材", "配伍/禁忌网络", "NetworkX + 力导向"],
]
add_table(s, 0.6, 1.5, 12.1, 4.8, ["功能", "输入", "输出", "技术"], data,
          col_w=[2.6, 2.8, 3.6, 3.1], fsize=13)
footnote(s)

# ---------------- 12. 官方评估指标 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "九、官方评估指标", "实测（evaluation/reports/eval_official_report.md）")
tf = box(s, 0.6, 1.5, 12.1, 5.3); tf.text_frame.word_wrap = True
set_text(tf, "评估设定", size=17, color=INK, bold=True)
add_para(tf, "数据集：val 10000 张 / 163 类；权重 best_model.pth（Epoch 5，510.9 MB）。", size=14, bullet=True)
add_para(tf, "推理：有文本走多模态，无文本走纯视觉分支。", size=14, bullet=True)
set_text(tf, "核心指标", size=17, color=INK, bold=True)
add_para(tf, "多模态（有文本）：Accuracy 0.9995，Top-5 1.0000", size=16, color=VERM, bold=True, bullet=True)
add_para(tf, "纯视觉（无文本）：Accuracy 0.9548，Top-5 0.9965", size=16, color=VERM, bold=True, bullet=True)
add_para(tf, "文本辅助增益 Δacc：+0.0447", size=16, bullet=True)
set_text(tf, "错误分析（诚实呈现）", size=17, color=INK, bold=True)
add_para(tf, "163 类中 159 类 precision/recall/f1 全部 = 1.00；错误仅 4 类，且全是「同药不同形态」（首乌藤块/片、天麻块/片）互混——人为拆类的固有难度。", size=14, bullet=True)
add_para(tf, "最难的首乌藤块/片模型仍能 95% 区分。", size=14, color=VERM, bullet=True)
footnote(s)

# ---------------- 13. 技术栈与部署 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "十、技术栈与部署", "8G 显存友好，纯 Python 服务无需前端构建")
lt = box(s, 0.6, 1.5, 6.0, 5.2); lt.text_frame.word_wrap = True
set_text(lt, "技术栈", size=17, color=INK, bold=True)
add_para(lt, "视觉：timm（Swin-Tiny 等）", size=14, bullet=True)
add_para(lt, "文本：Transformers（BERT）", size=14, bullet=True)
add_para(lt, "融合/分类：PyTorch", size=14, bullet=True)
add_para(lt, "图谱：NetworkX（CSV 驱动）", size=14, bullet=True)
add_para(lt, "检索：本地 BERT 语义向量（CPU）", size=14, bullet=True)
add_para(lt, "服务：FastAPI + Gradio；前端原生 HTML/CSS/JS", size=14, bullet=True)
rt = box(s, 6.9, 1.5, 6.0, 5.2); rt.text_frame.word_wrap = True
set_text(rt, "部署与硬件", size=17, color=INK, bold=True)
add_para(rt, "面向 8G 显存（bs=16, 224；OOM 降至 8 / 192）。", size=14, bullet=True)
add_para(rt, "一键环境：setup_env.bat / .sh；自检 check_environment.py。", size=14, bullet=True)
add_para(rt, "启动：python main.py --mode serve（:8000）/ --mode demo（:7862）。", size=14, bullet=True)
add_para(rt, "默认绑定 0.0.0.0，同局域网可直访；公网借内网穿透/云服务器。", size=14, bullet=True)
add_para(rt, "端侧：TorchScript 导出纯视觉分支，快速识图。", size=14, color=VERM, bullet=True)
footnote(s)

# ---------------- 14. 创新点 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "十一、创新点", "核心 / 体验 / 安全 三层创新")
tf = box(s, 0.6, 1.5, 12.1, 5.3); tf.text_frame.word_wrap = True
inn = [
    "多模态互补识别：视觉+文本+HCA 融合，文本辅助增益 +4.47 点，模拟中医「望闻问切」互补。",
    "双分支保底架构：纯视觉分支保证无文本也能识别，消除「缺文本即失效」单点风险。",
    "图谱驱动解析与推荐：识别自动联动药性/归经/功效/毒性，并基于十八反/十九畏做禁忌规避，形成闭环。",
    "RAG 可溯源问答：本地语义检索 + 来源标注，离线可用、自动降级，显著降低幻觉。",
    "可解释与安全并重：Grad-CAM + 推理路径 + 毒性强制警示 + 置信度提示，做成默认能力。",
    "轻量可落地：8G 显存友好、纯前端零构建、REST 可服务化、端侧可导出。",
]
for i, t in enumerate(inn):
    add_para(tf, f"{i+1}. {t}", size=14.5, color=DARK, bold=(i % 2 == 0), bullet=False, space=7)
footnote(s)

# ---------------- 15. 不足与改进方向（双栏） ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "十二、不足与改进方向", "诚实自检：把短板讲成可信度")
sec(s, 0.6, 1.45, 6.0, 0.5, "当前不足", VERM)
lt = box(s, 0.6, 2.05, 6.0, 4.8); lt.text_frame.word_wrap = True
short = [
    "数据规模有限：识别 163 类/图谱 200+，距万级品类差距大；长尾靠小样本补足。",
    "块/片形态易混：同药不同形态（首乌藤、天麻）仍互混，需细粒度标注。",
    "图谱偏静态：内存版规模受限，方剂仅 12 首，深度与时效性不足。",
    "文本依赖人工标注：未做古籍/药典自动抽取，RAG 语料覆盖有限。",
    "LLM 受限于外部 Key：无 Key 退化为本地检索，解释性下降。",
    "缺量化与对抗评测：光照/角度/跨域鲁棒性未系统测试。",
    "安全边界有限：无在线药师审核、个性化风险评估。",
]
for t in short:
    add_para(lt, t, size=12.5, color=DARK, bullet=True, space=5)
sec(s, 6.9, 1.45, 6.0, 0.5, "改进方向")
rt = box(s, 6.9, 2.05, 6.0, 4.8); rt.text_frame.word_wrap = True
imp = [
    "扩充样本与细粒度形态标注（含块/片变体）。",
    "接入 Neo4j 持久化图谱，支持更大规模推理。",
    "引入自动知识抽取（LLM/规则）丰富图谱与语料。",
    "增加古籍/药典自动抽取，扩展 RAG 覆盖。",
    "接入更稳的 LLM 或本地小模型，提升无 Key 体验。",
    "建立鲁棒性与对抗评测基准（光照/角度/跨域）。",
    "构建「识别+风险分级+人工复核」闭环，明确科普/诊断边界。",
]
for t in imp:
    add_para(rt, t, size=12.5, color=INK, bullet=True, space=5)
footnote(s)

# ---------------- 16. 演示截图墙 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "十三、系统演示截图", "新中式 Web 前端五大页签")
imgs = ["shihu.jpg", "honghua.jpg", "jinyihua.jpg", "chenpi.jpg"]
pos = [(0.5, 1.5), (6.9, 1.5), (0.5, 4.4), (6.9, 4.4)]
for fn, (l, t) in zip(imgs, pos):
    p = show(fn)
    if p:
        pic_fit(s, p, l, t, 6.0, 2.7)
footnote(s)

# ---------------- 17. 总结与交付物 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s)
title_bar(s, "十四、总结与交付物", "Summary & Deliverables")
lt = box(s, 0.6, 1.5, 6.0, 5.2); lt.text_frame.word_wrap = True
set_text(lt, "总结", size=17, color=INK, bold=True)
add_para(lt, "一套「识别—解析—推荐—可解释」闭环的中草药多模态智能体，8G 显存下达多模态 0.9995 / 纯视觉 0.9548。", size=13.5, bullet=True)
add_para(lt, "亮点：跨模态融合、图谱驱动推荐、RAG 可溯源问答、安全可解释；不足集中在数据规模、形态细粒度与知识深度，已列改进路径。", size=13.5, bullet=True)
rt = box(s, 6.9, 1.5, 6.0, 5.2); rt.text_frame.word_wrap = True
set_text(rt, "交付物（docs/交付文档/）", size=17, color=INK, bold=True)
add_para(rt, "项目介绍.pptx（本演示）", size=13.5, bullet=True)
add_para(rt, "用户手册 / 部署说明", size=13.5, bullet=True)
add_para(rt, "评估报告（evaluation/reports/）", size=13.5, bullet=True)
add_para(rt, "一键环境脚本与端侧导出工具", size=13.5, bullet=True)
add_para(rt, "郑重声明：科普辅助工具，结果仅供参考，严禁自行开方用药。", size=13.5, color=VERM, bold=True, bullet=True)
footnote(s)

# ---------------- 封底 ----------------
s = prs.slides.add_slide(BLANK); add_bg(s, INK)
tf = box(s, 1.0, 2.8, 11.3, 2.0); tf.text_frame.word_wrap = True
set_text(tf, "本草掠影 · 让识药更可信", size=34, color=RICE, bold=True, align=PP_ALIGN.CENTER)
add_para(tf, "Thanks — 中草药多模态识别系统", size=16, color=GOLD, align=PP_ALIGN.CENTER, space=0)
tf2 = box(s, 1.0, 6.3, 11.3, 0.6)
set_text(tf2, "免责声明：本系统为中医药科普辅助工具，不构成医疗建议。", size=12, color=RICE, align=PP_ALIGN.CENTER)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
final = OUT
try:
    prs.save(OUT)
except PermissionError:
    # 目标文件可能被 PowerPoint 打开而锁定，回退到另一文件名避免丢失成果
    final = OUT.replace(".pptx", "_新版.pptx")
    prs.save(final)
    print("WARN: 原文件被占用，已保存到:", final)
print("saved:", final, "slides:", len(prs.slides._sldIdLst))
print("verify exists:", os.path.isfile(final), "size:", os.path.getsize(final) if os.path.isfile(final) else -1)
print("cwd:", os.getcwd())
